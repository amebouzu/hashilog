"""走ログ 事業提案書 PPTX 生成スクリプト (洗練版).

docs/PROPOSAL.md の内容を 11 枚のピッチデックに圧縮し、よりブラッシュアップした
タイポグラフィ・余白・アクセント装飾で描画する。

設計コンセプト:
- 16:9
- 日本語: Yu Gothic UI (UI Bold は別途指定)
- カバー: 左に縦の赤アクセント、ロゴ・タイトル・タグラインの3層
- 各コンテンツスライド: タイトル左に細い赤ライン、右上にロゴ、右下にスライド番号
- カラム見出し: 番号付きの赤バッジ + 太字
- 表: 1行目がレッドアンダーラインのヘッダ、偶数行を薄グレーでゼブラ
- フッター: 小さく中央 (セクションタイトル時のみ表示)
- 大きな数字スライド: サービスの定量的ハイライトを 1 ページに集約

実行:
    python scripts/build_proposal_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# Constants & palette
# ---------------------------------------------------------------------------
JP_FONT = "Yu Gothic UI"
ROOT = Path(__file__).resolve().parents[1]
LOGO = ROOT / "public" / "logo.png"

RED = RGBColor(0xE1, 0x06, 0x00)
RED_DARK = RGBColor(0xB2, 0x04, 0x00)
INK = RGBColor(0x18, 0x18, 0x1B)
GRAY_900 = RGBColor(0x27, 0x27, 0x2A)
GRAY_700 = RGBColor(0x3F, 0x3F, 0x46)
GRAY_500 = RGBColor(0x71, 0x71, 0x7A)
GRAY_300 = RGBColor(0xD4, 0xD4, 0xD8)
GRAY_200 = RGBColor(0xE4, 0xE4, 0xE7)
GRAY_100 = RGBColor(0xF4, 0xF4, 0xF5)
GRAY_50 = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xFC, 0xFC, 0xFD)


# Slide dimensions (16:9)
SW = Inches(13.333)
SH = Inches(7.5)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, left, top, width, height, *, fill=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.fill.background()
    return box.text_frame


def set_run(run, text, *, size=14, bold=False, color=INK, font=JP_FONT, italic=False):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def set_para(tf, text, *, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
             clear=True, line_spacing=None, space_after=None, font=JP_FONT):
    if clear and tf.paragraphs and tf.paragraphs[0].text == "":
        para = tf.paragraphs[0]
    else:
        para = tf.add_paragraph()
    para.alignment = align
    if line_spacing is not None:
        para.line_spacing = line_spacing
    if space_after is not None:
        para.space_after = Pt(space_after)
    run = para.add_run()
    set_run(run, text, size=size, bold=bold, color=color, font=font)
    return para


def add_bullets(tf, items, *, size=12, color=INK, marker="■", marker_color=RED):
    """マーカー付きの箇条書き。マーカーと本文で色を切り替えられる。"""
    first_done = False
    for line in items:
        if not first_done and tf.paragraphs and tf.paragraphs[0].text == "":
            para = tf.paragraphs[0]
            first_done = True
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = 1.25
        para.space_after = Pt(2)
        run_m = para.add_run()
        set_run(run_m, marker + " ", size=size - 2, bold=True, color=marker_color)
        run_t = para.add_run()
        set_run(run_t, line, size=size, color=color)


def add_rect(slide, left, top, width, height, *, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    return sh


# ---------------------------------------------------------------------------
# Page chrome
# ---------------------------------------------------------------------------
def add_page_chrome(slide, *, page_no=None, total=None):
    """共通の右上ロゴ + 右下ページ番号 + フッター情報。"""
    # Logo top-right
    if LOGO.exists():
        h = Inches(0.42)
        w = Inches(0.42 * (1181 / 625))
        slide.shapes.add_picture(
            str(LOGO),
            SW - w - Inches(0.45),
            Inches(0.35),
            width=w,
            height=h,
        )

    # Footer line
    add_rect(
        slide,
        Inches(0.6),
        SH - Inches(0.45),
        SW - Inches(1.2),
        Emu(6350),  # ~0.5pt
        fill=GRAY_200,
    )

    # Footer text (left)
    tf = add_textbox(slide, Inches(0.6), SH - Inches(0.4), Inches(8), Inches(0.3))
    set_para(
        tf,
        "走ログ (Hashilog)   /   RBS   /   hashilog2024@gmail.com",
        size=8,
        color=GRAY_500,
    )

    # Page number (right)
    if page_no is not None:
        tf = add_textbox(
            slide, SW - Inches(2.6), SH - Inches(0.4), Inches(2.0), Inches(0.3)
        )
        text = f"{page_no:02d} / {total:02d}" if total else f"{page_no:02d}"
        set_para(tf, text, size=9, color=GRAY_500, align=PP_ALIGN.RIGHT)


def add_section_header(slide, num, title):
    """セクションヘッダ: 左に番号バッジ、太字タイトル、下に細い赤アクセント線。"""
    # Number badge (red square)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6),
        Inches(0.6),
        Inches(0.55),
        Inches(0.55),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RED
    badge.line.fill.background()
    btf = badge.text_frame
    btf.margin_left = Inches(0.0)
    btf.margin_top = Inches(0.05)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = btf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, str(num), size=20, bold=True, color=WHITE)

    # Title
    tf = add_textbox(slide, Inches(1.35), Inches(0.55), Inches(11), Inches(0.7))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_para(tf, title, size=24, bold=True, color=INK)

    # Accent underline
    add_rect(
        slide,
        Inches(0.6),
        Inches(1.32),
        Inches(0.8),
        Inches(0.04),
        fill=RED,
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_cover(prs):
    s = slide_blank(prs)

    # Off-white background tint
    bg = add_rect(s, 0, 0, SW, SH, fill=OFFWHITE)

    # Diagonal red accent on the left edge (vertical band)
    add_rect(s, 0, 0, Inches(0.6), SH, fill=RED)

    # Tiny red square accent
    add_rect(
        s,
        Inches(1.1),
        Inches(1.0),
        Inches(0.18),
        Inches(0.18),
        fill=RED,
    )

    # Eyebrow text
    tf = add_textbox(s, Inches(1.4), Inches(0.95), Inches(8), Inches(0.4))
    set_para(
        tf,
        "BUSINESS PROPOSAL",
        size=11,
        bold=True,
        color=RED,
    )

    # Logo (中央左寄せ大きめ)
    if LOGO.exists():
        h = Inches(2.2)
        w = Inches(2.2 * (1181 / 625))
        s.shapes.add_picture(str(LOGO), Inches(1.1), Inches(1.6), width=w, height=h)

    # Title (大判)
    tf = add_textbox(s, Inches(1.1), Inches(4.1), Inches(11), Inches(1.3))
    set_para(tf, "事業提案書", size=56, bold=True, color=INK)

    # Tagline
    tf = add_textbox(s, Inches(1.1), Inches(5.4), Inches(11), Inches(0.6))
    set_para(
        tf,
        "日本のサーキットタイムを、もっと面白く。もっと信頼できるものに。",
        size=18,
        color=GRAY_700,
    )

    # Bottom info line
    add_rect(
        s,
        Inches(1.1),
        Inches(6.45),
        Inches(0.5),
        Inches(0.04),
        fill=RED,
    )
    tf = add_textbox(s, Inches(1.1), Inches(6.55), Inches(11), Inches(0.6))
    set_para(
        tf,
        "RBS    /    久米田 昴    /    https://hashilog.jp    /    hashilog2024@gmail.com",
        size=11,
        color=GRAY_500,
    )


def slide_section_divider(prs, num, title, subtitle, *, page_no, total):
    """章扉スライド: 大きな番号 + タイトル + サブタイトル。"""
    s = slide_blank(prs)
    add_rect(s, 0, 0, SW, SH, fill=GRAY_50)
    add_rect(s, 0, 0, Inches(0.6), SH, fill=RED)

    # Big number on left
    tf = add_textbox(s, Inches(1.1), Inches(2.0), Inches(4.0), Inches(2.6))
    set_para(
        tf,
        f"{num:02d}",
        size=140,
        bold=True,
        color=RED,
    )

    # Title
    tf = add_textbox(s, Inches(5.0), Inches(2.7), Inches(8), Inches(1.0))
    set_para(tf, title, size=36, bold=True, color=INK)

    # Subtitle
    tf = add_textbox(s, Inches(5.0), Inches(4.0), Inches(8), Inches(1.5))
    set_para(tf, subtitle, size=15, color=GRAY_700, line_spacing=1.5)

    add_page_chrome(s, page_no=page_no, total=total)


def slide_two_col(prs, num, title, lead, columns, *, page_no, total):
    """2 カラム式の汎用スライド。各カラムは番号バッジ付き見出し + 箇条書き。"""
    s = slide_blank(prs)
    add_section_header(s, num, title)

    if lead:
        tf = add_textbox(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.7))
        set_para(tf, lead, size=13, color=GRAY_700)
        top = Inches(2.3)
    else:
        top = Inches(1.7)

    col_w = Inches(5.95)
    col_gap = Inches(0.4)
    for i, (heading, items) in enumerate(columns[:2]):
        left = Inches(0.6) + i * (col_w + col_gap)

        # Card-like background for the column
        add_rect(s, left, top, col_w, Inches(4.6), fill=GRAY_50)
        add_rect(s, left, top, Inches(0.06), Inches(4.6), fill=RED)

        # Column heading
        tf = add_textbox(
            s, left + Inches(0.3), top + Inches(0.2), col_w - Inches(0.4), Inches(0.55)
        )
        set_para(tf, heading, size=16, bold=True, color=INK)

        # Bullets
        tf = add_textbox(
            s,
            left + Inches(0.3),
            top + Inches(0.85),
            col_w - Inches(0.4),
            Inches(3.5),
        )
        add_bullets(tf, items, size=12)

    add_page_chrome(s, page_no=page_no, total=total)


def slide_three_col(prs, num, title, lead, columns, *, page_no, total):
    s = slide_blank(prs)
    add_section_header(s, num, title)

    if lead:
        tf = add_textbox(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.7))
        set_para(tf, lead, size=13, color=GRAY_700)
        top = Inches(2.3)
    else:
        top = Inches(1.7)

    gap = Inches(0.25)
    total_w = Inches(12.13)
    col_w = (total_w - gap * 2) / 3

    for i, (heading, items) in enumerate(columns[:3]):
        left = Inches(0.6) + i * (col_w + gap)

        add_rect(s, left, top, col_w, Inches(4.6), fill=GRAY_50)
        add_rect(s, left, top, Inches(0.06), Inches(4.6), fill=RED)

        tf = add_textbox(
            s, left + Inches(0.25), top + Inches(0.2), col_w - Inches(0.35), Inches(0.55)
        )
        set_para(tf, heading, size=14, bold=True, color=INK)

        tf = add_textbox(
            s,
            left + Inches(0.25),
            top + Inches(0.8),
            col_w - Inches(0.35),
            Inches(3.6),
        )
        add_bullets(tf, items, size=10.5)

    add_page_chrome(s, page_no=page_no, total=total)


def slide_table(prs, num, title, lead, headers, rows, col_widths_emu=None,
                *, page_no, total, first_col_bold=False):
    s = slide_blank(prs)
    add_section_header(s, num, title)

    if lead:
        tf = add_textbox(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.7))
        set_para(tf, lead, size=13, color=GRAY_700)
        top = Inches(2.3)
    else:
        top = Inches(1.7)

    cols = len(headers)
    rows_count = len(rows) + 1
    table_height = min(Inches(0.5) + Inches(0.42) * len(rows), Inches(4.7))
    table_shape = s.shapes.add_table(
        rows_count, cols, Inches(0.6), top, Inches(12.13), table_height
    )
    table = table_shape.table

    if col_widths_emu:
        for i, w in enumerate(col_widths_emu):
            table.columns[i].width = w

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        set_run(r, h, size=11, bold=True, color=WHITE)

    # Body rows (zebra)
    for ri, row in enumerate(rows, start=1):
        bg = WHITE if ri % 2 == 1 else GRAY_50
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            is_first = ci == 0
            set_run(
                r,
                val,
                size=10.5,
                color=INK,
                bold=first_col_bold and is_first,
            )

    add_page_chrome(s, page_no=page_no, total=total)


def slide_stats(prs, num, title, lead, stats, *, page_no, total):
    """大きな数字を 4 つ並べる定量的ハイライトスライド."""
    s = slide_blank(prs)
    add_section_header(s, num, title)

    if lead:
        tf = add_textbox(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.7))
        set_para(tf, lead, size=13, color=GRAY_700)

    top = Inches(2.5)
    gap = Inches(0.25)
    n = len(stats)
    total_w = Inches(12.13)
    card_w = (total_w - gap * (n - 1)) / n
    card_h = Inches(3.2)

    for i, (number, label, sub) in enumerate(stats):
        left = Inches(0.6) + i * (card_w + gap)
        # Card
        add_rect(s, left, top, card_w, card_h, fill=GRAY_50)
        add_rect(s, left, top, card_w, Inches(0.08), fill=RED)

        # Big number
        tf = add_textbox(s, left, top + Inches(0.4), card_w, Inches(1.6))
        set_para(
            tf,
            number,
            size=64,
            bold=True,
            color=RED,
            align=PP_ALIGN.CENTER,
        )

        # Label
        tf = add_textbox(s, left, top + Inches(2.05), card_w, Inches(0.6))
        set_para(
            tf,
            label,
            size=14,
            bold=True,
            color=INK,
            align=PP_ALIGN.CENTER,
        )

        # Sub
        tf = add_textbox(s, left + Inches(0.2), top + Inches(2.55), card_w - Inches(0.4), Inches(0.7))
        set_para(
            tf,
            sub,
            size=10,
            color=GRAY_500,
            align=PP_ALIGN.CENTER,
            line_spacing=1.3,
        )

    add_page_chrome(s, page_no=page_no, total=total)


def slide_closing(prs, *, page_no, total):
    s = slide_blank(prs)
    # Dark hero background top half
    add_rect(s, 0, 0, SW, Inches(4.5), fill=INK)
    # Bottom contact area light
    add_rect(s, 0, Inches(4.5), SW, SH - Inches(4.5), fill=OFFWHITE)
    # Red accent strip
    add_rect(s, 0, Inches(4.42), SW, Inches(0.08), fill=RED)

    # Logo (white area would clash on dark, so omit on dark; place small on bottom)
    add_page_chrome(s, page_no=page_no, total=total)

    # Hero copy
    tf = add_textbox(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.2))
    set_para(
        tf,
        "サーキットを走るすべての人が、",
        size=34,
        bold=True,
        color=WHITE,
    )
    tf = add_textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2))
    set_para(
        tf,
        "自分のタイムをもっと誇れる場所へ。",
        size=34,
        bold=True,
        color=RED,
    )

    tf = add_textbox(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(1.0))
    set_para(
        tf,
        "ぜひ一度、サービスを実際に触っていただき、提携の可能性についてお話しさせてください。",
        size=14,
        color=GRAY_300,
    )

    # Contact card
    card_top = Inches(4.85)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.9),
        card_top,
        Inches(11.5),
        Inches(2.0),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = GRAY_200
    card.shadow.inherit = False

    # 中身
    tf = add_textbox(
        s,
        Inches(1.2),
        card_top + Inches(0.2),
        Inches(11.0),
        Inches(0.6),
    )
    set_para(tf, "走ログ (Hashilog)", size=20, bold=True, color=INK)

    info_top = card_top + Inches(0.85)
    info = [
        ("Web", "https://hashilog.jp"),
        ("事業者", "RBS"),
        ("責任者", "久米田 昴"),
        ("お問い合わせ", "hashilog2024@gmail.com"),
    ]
    for i, (k, v) in enumerate(info):
        col = i % 2
        row = i // 2
        left = Inches(1.2) + col * Inches(5.2)
        top = info_top + row * Inches(0.45)
        tf = add_textbox(s, left, top, Inches(5.0), Inches(0.4))
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        rk = para.add_run()
        set_run(rk, f"{k}", size=10, bold=True, color=GRAY_500)
        rs = para.add_run()
        set_run(rs, "    " + v, size=12, color=INK)


# ---------------------------------------------------------------------------
# Main build
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # Total content slides (excluding cover and closing for the chrome counter)
    # We'll just count actual slides in the deck.
    TOTAL = 11

    # 1. Cover
    slide_cover(prs)

    # 2. Service overview / value
    slide_two_col(
        prs,
        1,
        "サービス概要・提供価値",
        "サーキットで刻んだタイムを、愛車情報と一緒にシェアできる Web サービス。本番稼働中。",
        [
            (
                "ドライバーに",
                [
                    "タイム履歴・改造履歴・タイヤ履歴を1か所で管理",
                    "X / Insta / Threads / FB / LINE で 1-click シェア",
                    "同条件ドライバーのセッティングを参考にできる",
                    "前後別タイヤ・気温・路温まで記録できる細かさ",
                    "全国 25+ サーキットを横断比較",
                ],
            ),
            (
                "事業者に",
                [
                    "サーキット運営: 公式ページ + 走行会告知 (無料)",
                    "タイヤ・パーツメーカー: 自社銘柄の使用実績データ",
                    "イベント主催者: 走行会の集客チャネル",
                    "媒体: 走行データの記事引用・共同企画",
                    "自動車メーカー: 車種別のサーキット適性指標",
                ],
            ),
        ],
        page_no=2,
        total=TOTAL,
    )

    # 3. Market & approach
    slide_two_col(
        prs,
        2,
        "市場背景と解決アプローチ",
        "国内 30+ サーキット、参加層は数万人。30〜50 代男性中心で可処分所得が高い。",
        [
            (
                "現状の課題",
                [
                    "個人ブログ・SNS のスクショ・LINE グループに散在",
                    "サーキット間で形式バラバラ、車両/タイヤと紐づかない",
                    "走行会の Excel は主催者依存、横断比較できない",
                    "過去データに辿り着けない、検索性ゼロ",
                ],
            ),
            (
                "走ログのアプローチ",
                [
                    "検索可能な構造化フォーマット (多次元フィルター)",
                    "エビデンス文化 (計測機器・モニター写真の添付)",
                    "SNS 連携で OGP 最適化、外部拡散を前提",
                    "サーキット運営者の公式ページを横断的に統合",
                ],
            ),
        ],
        page_no=3,
        total=TOTAL,
    )

    # 4. Stats highlight (新スライド: ブラッシュアップ要素)
    slide_stats(
        prs,
        3,
        "数字で見る走ログ",
        "プロダクトのスケールと粒度を一目で。",
        [
            ("25+", "サーキット登録", "全国の国際格式・ミニ・ショート\nコースを網羅"),
            ("6+", "SNS 連携", "X / Insta / Threads / FB\nLINE / リンクコピー"),
            ("2系統", "前後別タイヤ", "ブランド・銘柄・サイズを\n前後それぞれ記録可能"),
            ("100%", "実装済み", "MVP 機能はすべて本番稼働\n継続的にデプロイ中"),
        ],
        page_no=4,
        total=TOTAL,
    )

    # 5. Implementation (3 col)
    slide_three_col(
        prs,
        4,
        "実装済み機能ハイライト",
        "「設計だけ」ではなく、すでに本番稼働しているプロダクト (2026年5月時点)。",
        [
            (
                "ユーザー / 愛車",
                [
                    "メール+パスワード認証、退会で完全削除",
                    "プロフィール: 自己紹介・都道府県・SNS 6 種",
                    "マイガレージ: 複数台、改造内容を分野別に",
                    "写真ギャラリー (カバー + 複数枚)",
                ],
            ),
            (
                "タイム投稿 (中核)",
                [
                    "総合タイム、セクター 1〜4、最高速",
                    "天候・路面・気温・路温",
                    "前後別タイヤ銘柄・サイズ (差別化)",
                    "ユーザー入力銘柄は自動登録",
                    "エビデンス写真複数枚、編集・削除可",
                ],
            ),
            (
                "ランキング / B2B / 信頼性",
                [
                    "サーキット × 車種 × タイヤで多次元フィルター",
                    "サーキット運営者の公式アカウント・告知機能",
                    "OGP/JSON-LD/PWA/sitemap 完備",
                    "Vercel + Supabase Tokyo (国内データ保管)",
                    "通報・モデレーション、規約・特商法整備",
                ],
            ),
        ],
        page_no=5,
        total=TOTAL,
    )

    # 6. Target & B2B
    slide_two_col(
        prs,
        5,
        "ターゲット",
        "ドライバーは可処分所得が高くロイヤリティが高い。B2B は複数業界に展開可能。",
        [
            (
                "メインユーザー (ドライバー)",
                [
                    "タイムアタッカー: 30〜50 代、年 10〜30 走行日",
                    "走行会レギュラー: 月 1〜2 回、エビデンス志向",
                    "若手 SNS ネイティブ: 20 代、発信意欲旺盛",
                    "耐久・チーム参戦: 同一車両で複数ドライバー比較",
                ],
            ),
            (
                "B2B パートナー候補",
                [
                    "サーキット運営会社 (公式ページ+告知)",
                    "タイヤ・パーツメーカー (使用実績データ)",
                    "自動車メーカー (車種別サーキット適性)",
                    "走行会主催者・モータースポーツ媒体",
                    "損保・ファイナンス (ターゲティング広告)",
                ],
            ),
        ],
        page_no=6,
        total=TOTAL,
    )

    # 7. Monetization (table)
    slide_table(
        prs,
        6,
        "収益化モデル",
        "コア機能は永久無料を維持。Stripe 課金基盤実装済 (フラグでON/OFF)。",
        ["対象", "プラン / メニュー", "内容"],
        [
            ["B2C", "Free  (0円)", "全コア機能 (タイム投稿・閲覧・シェア・ランキング)"],
            ["B2C", "Premium  (480円/月程度)", "写真大量UL、推移グラフ、CSV、広告非表示"],
            ["B2C", "Pro  (980円/月程度)", "チームアカウント、複数車両比較、コーチング"],
            ["B2B サーキット", "公式運営者 (無料)", "自施設ページ編集・走行会告知"],
            ["B2B サーキット", "プロモーション (要相談)", "トップ露出・ニュースレター・送客レポート"],
            ["B2B メーカー", "公式バッジ / データ提供", "ブランド表示優先・月次集計レポート"],
            ["B2B メーカー", "タイアップ・広告枠・スポンサー", "新製品特集・バナー・走ログ大会協賛"],
        ],
        col_widths_emu=[Inches(2.0), Inches(3.6), Inches(6.5)],
        first_col_bold=True,
        page_no=7,
        total=TOTAL,
    )

    # 8. Differentiation (table)
    slide_table(
        prs,
        7,
        "競合・差別化",
        "既存の選択肢と比較した走ログのポジション。",
        ["項目", "ブログ", "SNS", "サーキット公式", "走ログ"],
        [
            ["検索性", "低", "中", "中", "高 (多次元)"],
            ["データ構造化", "×", "×", "△", "○"],
            ["エビデンス", "△", "△", "○", "○"],
            ["横断比較", "×", "×", "×", "○"],
            ["公式情報統合", "×", "×", "○ (自社のみ)", "○ (横断)"],
            ["SNS シェア最適化", "△", "○", "×", "○"],
            ["改造内容と紐付け", "△", "△", "×", "○"],
            ["タイヤ前後別記録", "×", "×", "×", "○"],
        ],
        first_col_bold=True,
        page_no=8,
        total=TOTAL,
    )

    # 9. Partnership (3 col)
    slide_three_col(
        prs,
        8,
        "パートナーシップ提案",
        "業種別にご提供できる内容。詳細は個別ご相談ください。",
        [
            (
                "サーキット運営者",
                [
                    "公式ページ編集 (説明・代表コーナー・URL)",
                    "走行会・レース・お知らせ告知",
                    "自施設のラップタイム集計",
                    "スタッフアカウント発行",
                    "→ 必要なのはご担当のメアドのみ。無料。",
                ],
            ),
            (
                "タイヤ・パーツメーカー",
                [
                    "使用実績データ (匿名化集計、月次)",
                    "ブランド公式バッジ表示",
                    "新製品リリース連動の特集枠",
                    "コラボ企画 (タイムアタック / プレゼント)",
                    "サンプリング協力 (モニター募集)",
                ],
            ),
            (
                "走行会主催者・媒体",
                [
                    "イベント告知ページ (無料)",
                    "参加者のタイム集計テンプレート",
                    "イベント特化ページ作成 (要相談)",
                    "走ログ内データの記事引用 (媒体)",
                    "共同企画記事・アンケート協力",
                ],
            ),
        ],
        page_no=9,
        total=TOTAL,
    )

    # 10. Roadmap (2 col)
    slide_two_col(
        prs,
        9,
        "ロードマップ",
        "MVP を本番稼働させた上で、データ価値とパートナー連携を深掘りしていく。",
        [
            (
                "完了済み (2026年5月)",
                [
                    "認証 / プロフィール / SNS 連携",
                    "愛車登録・改造履歴",
                    "タイム投稿 (前後別タイヤ含む)",
                    "ランキング・多次元フィルター",
                    "サーキット 25+ 登録 / 運営者管理",
                    "通報・モデレーション",
                    "PWA / SEO / 構造化データ",
                    "規約・プライバシー・特商法整備",
                ],
            ),
            (
                "短期〜中期 (3〜12 か月)",
                [
                    "ベストタイム推移グラフ",
                    "同一車種比較ダッシュボード",
                    "走行会・イベントカレンダー (横断)",
                    "通知機能 (記録更新・フォロー)",
                    "メーカー向け管理コンソール",
                    "チームアカウント / API 公開",
                    "動画オンボードアップロード",
                    "AI 走行ライン分析 (将来)",
                ],
            ),
        ],
        page_no=10,
        total=TOTAL,
    )

    # 11. Closing
    slide_closing(prs, page_no=11, total=TOTAL)

    out = ROOT / "docs" / "PROPOSAL.pptx"
    # PowerPoint が開いていると上書き失敗する。tmp を経由して可能ならリネーム。
    tmp = out.with_suffix(".pptx.tmp")
    prs.save(str(tmp))
    try:
        if out.exists():
            out.unlink()
        tmp.rename(out)
        print(f"Wrote {out} ({len(prs.slides)} slides)")
    except PermissionError:
        print(
            f"Wrote {tmp} ({len(prs.slides)} slides). "
            f"{out.name} is locked (PowerPoint で開いていませんか?). "
            f"閉じてから手動で {tmp.name} → {out.name} にリネームしてください。"
        )


if __name__ == "__main__":
    build()
